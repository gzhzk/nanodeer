"""One compact boundary for basic DOCX, XLSX, and PPTX artifacts."""

from __future__ import annotations

import os
import re
import tempfile
import zipfile
from contextlib import contextmanager
from pathlib import Path
from typing import Any, Iterator, Literal

from docx import Document
from langchain_core.tools import tool
from openpyxl import Workbook, load_workbook
from pptx import Presentation

from nanodeer.workspace import WorkspacePathError, resolve_workspace_path

_MAX_ARCHIVE_BYTES = 25 * 1024 * 1024
_MAX_UNCOMPRESSED_BYTES = 50 * 1024 * 1024
_MAX_INSPECT_CHARS = 20_000


@contextmanager
def _atomic_target(path: Path) -> Iterator[Path]:
    """Yield a same-directory temporary path and atomically publish it."""
    path.parent.mkdir(parents=True, exist_ok=True)
    fd, temporary = tempfile.mkstemp(
        prefix=f".{path.stem}.", suffix=path.suffix, dir=path.parent
    )
    os.close(fd)
    temp_path = Path(temporary)
    try:
        yield temp_path
        os.replace(temp_path, path)
    finally:
        temp_path.unlink(missing_ok=True)


def _check_archive(path: Path) -> None:
    """Reject oversized or suspicious OOXML packages before inspection."""
    if path.stat().st_size > _MAX_ARCHIVE_BYTES:
        raise ValueError("office artifact exceeds 25MB inspection limit")
    with zipfile.ZipFile(path) as archive:
        infos = archive.infolist()
        if len(infos) > 1000 or sum(info.file_size for info in infos) > _MAX_UNCOMPRESSED_BYTES:
            raise ValueError("office artifact expands beyond safe inspection limits")
        if any(name.startswith("/") or ".." in Path(name).parts for name in archive.namelist()):
            raise ValueError("office artifact contains unsafe archive paths")


def _create_docx(path: Path, title: str, content: str) -> None:
    document = Document()
    if title:
        document.add_heading(title, level=0)
    for line in content.splitlines() or ([content] if content else []):
        document.add_paragraph(line)
    with _atomic_target(path) as target:
        document.save(target)


def _sheet_title(title: str) -> str:
    cleaned = re.sub(r"[\\/*?:\[\]]", "-", title or "Sheet1").strip()
    return (cleaned or "Sheet1")[:31]


def _create_xlsx(path: Path, title: str, data: list[list[Any]]) -> None:
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = _sheet_title(title)
    for row in data:
        worksheet.append(row)
    with _atomic_target(path) as target:
        workbook.save(target)
    workbook.close()


def _create_pptx(
    path: Path,
    title: str,
    content: str,
    slides: list[dict[str, str]] | None,
) -> None:
    presentation = Presentation()
    slide_data = slides or [{"title": title or "Presentation", "body": content}]
    for index, item in enumerate(slide_data, start=1):
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = str(item.get("title") or f"Slide {index}")
        body = slide.placeholders[1]
        body.text = str(item.get("body") or "")
    with _atomic_target(path) as target:
        presentation.save(target)


def _inspect_docx(path: Path) -> str:
    document = Document(path)
    lines = [paragraph.text for paragraph in document.paragraphs if paragraph.text]
    for table in document.tables:
        lines.extend("\t".join(cell.text for cell in row.cells) for row in table.rows)
    return "\n".join(lines)


def _inspect_xlsx(path: Path) -> str:
    workbook = load_workbook(path, read_only=True, data_only=True)
    sections: list[str] = []
    try:
        for worksheet in workbook.worksheets:
            rows = [
                "\t".join("" if value is None else str(value) for value in row)
                for row in worksheet.iter_rows(values_only=True)
            ]
            sections.append(f"[Sheet: {worksheet.title}]\n" + "\n".join(rows))
    finally:
        workbook.close()
    return "\n\n".join(sections)


def _inspect_pptx(path: Path) -> str:
    presentation = Presentation(path)
    slides: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        lines = [shape.text for shape in slide.shapes if hasattr(shape, "text_frame") and shape.text]
        slides.append(f"[Slide {index}]\n" + "\n".join(lines))
    return "\n\n".join(slides)


def _inspect(path: Path) -> str:
    _check_archive(path)
    readers = {
        ".docx": _inspect_docx,
        ".xlsx": _inspect_xlsx,
        ".pptx": _inspect_pptx,
    }
    result = readers[path.suffix.lower()](path)
    return result[:_MAX_INSPECT_CHARS] or "(artifact contains no readable text)"


@tool
def office_artifact(
    action: Literal["create", "inspect"],
    file_path: str,
    title: str = "",
    content: str = "",
    data: list[list[Any]] | None = None,
    slides: list[dict[str, str]] | None = None,
) -> str:
    """Create or inspect a basic DOCX, XLSX, or PPTX artifact.

    Args:
        action: ``create`` or ``inspect``.
        file_path: Workspace path ending in .docx, .xlsx, or .pptx.
        title: Document title, worksheet name, or presentation title.
        content: Plain-text document body.
        data: Spreadsheet rows represented as a two-dimensional JSON array.
        slides: Presentation slides, each with ``title`` and ``body`` strings.
    """
    try:
        path = resolve_workspace_path(
            file_path, access="write" if action == "create" else "read"
        )
        suffix = path.suffix.lower()
        if suffix not in {".docx", ".xlsx", ".pptx"}:
            return "Error: file_path must end in .docx, .xlsx, or .pptx"
        if action == "inspect":
            if not path.is_file():
                return f"Error: office artifact not found: {file_path}"
            return _inspect(path)

        if suffix == ".docx":
            _create_docx(path, title, content)
        elif suffix == ".xlsx":
            if not data:
                return "Error: spreadsheet creation requires non-empty data"
            _create_xlsx(path, title, data)
        else:
            _create_pptx(path, title, content, slides)
        return f"Created {suffix[1:].upper()} artifact at {file_path} ({path.stat().st_size} bytes)"
    except (OSError, ValueError, KeyError, zipfile.BadZipFile, WorkspacePathError) as exc:
        return f"Error handling office artifact {file_path}: {exc}"


__all__ = ["office_artifact"]
